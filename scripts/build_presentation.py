"""
StayNep AI Voice Receptionist — Innovation Project Final Presentation
Generates a 18-slide .pptx matching the dark-theme format.
Run: python3 scripts/build_presentation.py
Output: StayNep_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
BG          = RGBColor(0x0B, 0x16, 0x1F)   # main background
BG_CARD     = RGBColor(0x0D, 0x22, 0x2E)   # card fill
BG_PANEL    = RGBColor(0x07, 0x12, 0x19)   # left panel on title slide
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEAL        = RGBColor(0x00, 0xD4, 0xAA)   # accent / label
LABEL_TEAL  = RGBColor(0x00, 0xD4, 0xAA)
MUTED       = RGBColor(0x8A, 0xA0, 0xB0)   # body / muted text
GREEN_BAR   = RGBColor(0x00, 0xD4, 0xAA)
BLUE_BAR    = RGBColor(0x4A, 0x9E, 0xFF)
YELLOW_BAR  = RGBColor(0xF5, 0xC5, 0x42)
PINK_BAR    = RGBColor(0xFF, 0x6B, 0x8A)
DIVIDER     = RGBColor(0x1A, 0x30, 0x3E)

W  = Inches(13.33)   # slide width  (widescreen 16:9)
H  = Inches(7.5)     # slide height

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill_color=BG_CARD, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def accent_bar(slide, x, y, h, color=GREEN_BAR, width=Inches(0.055)):
    b = box(slide, x, y, width, h, fill_color=color)
    return b


def txt(slide, text, x, y, w, h,
        size=Pt(14), bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tf = slide.shapes.add_textbox(x, y, w, h)
    tf.word_wrap = wrap
    p  = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tf


def label(slide, text, x, y, w=Inches(3)):
    """Small teal ALL-CAPS label above a heading."""
    txt(slide, text.upper(), x, y, w, Inches(0.3),
        size=Pt(8), bold=True, color=LABEL_TEAL)


def heading(slide, text, x, y, w, size=Pt(30)):
    txt(slide, text, x, y, w, Inches(0.9),
        size=size, bold=False, color=WHITE)


def subtext(slide, text, x, y, w):
    txt(slide, text, x, y, w, Inches(0.4),
        size=Pt(11), color=MUTED, italic=True)


def card(slide, title, body, x, y, w, h=Inches(1.5),
         bar_color=GREEN_BAR, title_size=Pt(13)):
    """Dark card with left accent bar."""
    box(slide, x, y, w, h, fill_color=BG_CARD)
    accent_bar(slide, x, y + Inches(0.18), h - Inches(0.36), color=bar_color)
    txt(slide, title, x + Inches(0.18), y + Inches(0.14), w - Inches(0.22),
        Inches(0.35), size=title_size, bold=True, color=WHITE)
    txt(slide, body, x + Inches(0.18), y + Inches(0.5), w - Inches(0.22),
        h - Inches(0.55), size=Pt(10), color=MUTED)


def stat(slide, value, label_text, x, y, val_color=TEAL):
    txt(slide, value, x, y, Inches(2), Inches(0.65),
        size=Pt(44), bold=False, color=val_color)
    txt(slide, label_text, x, y + Inches(0.6), Inches(2), Inches(0.4),
        size=Pt(10), color=MUTED)


def footer(slide, page, total=18):
    txt(slide, "StayNep AI Voice Receptionist", Inches(0.4), Inches(7.1),
        Inches(5), Inches(0.3), size=Pt(9), color=MUTED)
    txt(slide, f"{page:02d} / {total}", Inches(11.8), Inches(7.1),
        Inches(1.5), Inches(0.3), size=Pt(9), color=MUTED, align=PP_ALIGN.RIGHT)
    # bottom divider line
    line = slide.shapes.add_shape(1, Inches(0.4), Inches(6.98), Inches(12.53), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = DIVIDER
    line.line.fill.background()


def flow_step(slide, num, title, body, x, y, w=Inches(2.1), h=Inches(1.35)):
    box(slide, x, y, w, h, fill_color=BG_CARD)
    txt(slide, num, x + Inches(0.14), y + Inches(0.1), w, Inches(0.28),
        size=Pt(9), color=MUTED)
    txt(slide, title, x + Inches(0.14), y + Inches(0.32), w - Inches(0.2),
        Inches(0.35), size=Pt(13), bold=True, color=WHITE)
    txt(slide, body, x + Inches(0.14), y + Inches(0.65), w - Inches(0.2),
        h - Inches(0.7), size=Pt(9.5), color=MUTED)


def arrow(slide, x, y):
    """Small → arrow between flow steps."""
    txt(slide, "→", x, y, Inches(0.3), Inches(0.3),
        size=Pt(14), color=MUTED, align=PP_ALIGN.CENTER)


# ── SLIDES ────────────────────────────────────────────────────────────────────

def slide01_title(prs):
    s = blank_slide(prs)
    bg(s)
    # Left dark panel
    box(s, 0, 0, Inches(4.3), H, fill_color=BG_PANEL)
    # Logo placeholder
    txt(s, "⬡ NEP", Inches(0.35), Inches(0.45), Inches(1.5), Inches(0.5),
        size=Pt(14), bold=True, color=TEAL)
    # Main heading
    txt(s, "AI Voice Receptionist", Inches(0.35), Inches(1.3), Inches(3.7),
        Inches(0.85), size=Pt(30), bold=False, color=WHITE)
    txt(s, "for Hotels & Hospitality", Inches(0.35), Inches(2.1), Inches(3.7),
        Inches(0.5), size=Pt(22), color=TEAL)
    txt(s, "Intelligent 24/7 multilingual hotel reception\npowered by generative AI, live inventory, and\nautonomous booking.",
        Inches(0.35), Inches(2.75), Inches(3.7), Inches(0.9),
        size=Pt(11), color=MUTED)

    # Right side cards
    rx = Inches(4.7)
    card(s, "Group Name", "Team Think", rx, Inches(0.9), Inches(8.2),
         Inches(1.1), bar_color=YELLOW_BAR)
    card(s, "Team Members",
         "Prabin Sharma\nNikesh Kumar Sah\nDhrub Narayan Yadav\nRiyaj Khadka",
         rx, Inches(2.2), Inches(8.2), Inches(1.8), bar_color=BLUE_BAR)
    card(s, "Project", "StayNep AI Voice Receptionist",
         rx, Inches(4.2), Inches(8.2), Inches(1.1), bar_color=GREEN_BAR)

    txt(s, "Innovation Project Final Presentation",
        Inches(0.35), Inches(6.85), Inches(4), Inches(0.3),
        size=Pt(9), color=MUTED)


def slide02_exec_summary(prs):
    s = blank_slide(prs)
    bg(s)
    label(s, "Executive Summary", Inches(0.5), Inches(0.35))
    heading(s, "Hotels get an always-on multilingual front desk.",
            Inches(0.5), Inches(0.65), Inches(12.3), size=Pt(28))
    subtext(s, "What the project is, who it serves, and why it matters.",
            Inches(0.5), Inches(1.35), Inches(10))

    cw = Inches(5.9)
    ch = Inches(1.55)
    cy1, cy2 = Inches(1.75), Inches(3.45)

    card(s, "One-sentence pitch",
         "A voice-first AI receptionist that answers guest questions, checks availability, and completes bookings without front-desk handoff.",
         Inches(0.5), cy1, cw, ch, bar_color=GREEN_BAR)
    card(s, "Problem",
         "Hotels lose guest trust and bookings when staff are busy, offline, or unable to speak the guest's language.",
         Inches(6.6), cy1, cw, ch, bar_color=PINK_BAR)
    card(s, "Solution",
         "StayNep combines speech, Gemini/OpenAI reasoning, hotel-specific knowledge, live inventory, and confirmations in one assistant.",
         Inches(0.5), cy2, cw, ch, bar_color=BLUE_BAR)
    card(s, "Target users",
         "Hotels, resorts, guesthouses, homestays, and travel operators that need affordable 24/7 guest service.",
         Inches(6.6), cy2, cw, ch, bar_color=YELLOW_BAR)

    stat(s, "34", "supported\nlanguages", Inches(0.5), Inches(5.2), val_color=TEAL)
    stat(s, "24/7", "guest support", Inches(3.1), Inches(5.2), val_color=BLUE_BAR)
    stat(s, "0", "routine handoff\ntarget", Inches(6.2), Inches(5.2), val_color=YELLOW_BAR)
    footer(s, 2)


def slide03_problem(prs):
    s = blank_slide(prs)
    bg(s)
    label(s, "Problem Statement", Inches(0.5), Inches(0.35))
    heading(s, "Routine hotel service still depends on staff availability.",
            Inches(0.5), Inches(0.65), Inches(12.3), size=Pt(26))
    subtext(s, "The pain is most visible during off-hours, peak season, and multilingual guest interactions.",
            Inches(0.5), Inches(1.35), Inches(11))

    cw, ch = Inches(2.9), Inches(1.6)
    cx = [Inches(0.5), Inches(3.6), Inches(6.7), Inches(9.8)]
    cy = Inches(1.8)
    bars = [PINK_BAR, TEAL, YELLOW_BAR, BLUE_BAR]

    items = [
        ("Off-hours", "Guests still expect answers after reception desks slow down."),
        ("Language gaps", "International and domestic travelers may not share one service language."),
        ("Missed bookings", "Slow replies turn high-intent inquiries into lost revenue."),
        ("Staff load", "Front desks repeat the same FAQ and booking tasks every day."),
    ]
    for i, (title, body) in enumerate(items):
        card(s, title, body, cx[i], cy, cw, ch, bar_color=bars[i])

    txt(s, "Real-life example: a guest asks about room availability late at night, wants a reply in their own language, and expects confirmation immediately. A traditional front desk may need a callback; StayNep can continue the workflow in the same conversation.",
        Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.8),
        size=Pt(11), color=MUTED, italic=True)
    footer(s, 3)


def slide04_discovery(prs):
    s = blank_slide(prs)
    bg(s)
    label(s, "Customer Discovery & Empathy", Inches(0.5), Inches(0.35))
    heading(s, "The strongest signal is operational friction, not lack of websites.",
            Inches(0.5), Inches(0.65), Inches(12.3), size=Pt(25))
    subtext(s, "Discovery focused on hotel owners, reception teams, guests, and tourism businesses.",
            Inches(0.5), Inches(1.35), Inches(11))

    cw, ch = Inches(5.9), Inches(1.55)
    items = [
        ("Hotel owner",
         '"Hiring and retraining staff is expensive and inconsistent."\nNeed: lower routine service cost.',
         Inches(0.5), Inches(1.8), GREEN_BAR),
        ("Reception team",
         '"Peak-hour calls and repeated questions interrupt higher-value work."\nNeed: triage and automation.',
         Inches(6.6), Inches(1.8), BLUE_BAR),
        ("Guest",
         '"I want answers in my language without waiting for a person."\nNeed: instant, natural support.',
         Inches(0.5), Inches(3.55), YELLOW_BAR),
        ("Tourism operator",
         '"Slow inquiry follow-up can lose group bookings."\nNeed: faster conversion.',
         Inches(6.6), Inches(3.55), PINK_BAR),
    ]
    for title, body, x, y, bar in items:
        card(s, title, body, x, y, cw, ch, bar_color=bar)

    txt(s, "What we learned: guests do not want another static FAQ; hotels need an assistant that understands intent, checks real data, and knows when a human is truly required.",
        Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.7),
        size=Pt(11), color=MUTED, italic=True)
    footer(s, 4)


def slide05_solution(prs):
    s = blank_slide(prs)
    bg(s)
    label(s, "Proposed Solution", Inches(0.5), Inches(0.35))
    heading(s, "A complete guest service workflow runs inside the conversation.",
            Inches(0.5), Inches(0.65), Inches(12.3), size=Pt(26))
    subtext(s, "Voice input becomes action: answer, availability check, booking, dining reservation, cancellation, or escalation.",
            Inches(0.5), Inches(1.35), Inches(12))

    steps = [
        ("01", "Guest speaks", "Mic, phone, embed, or text input"),
        ("02", "STT", "Transcript is sanitized and validated"),
        ("03", "Intent router", "Booking, dining, FAQ, or escalation"),
        ("04", "Live data", "Prisma inventory and hotel config"),
        ("05", "Response", "Streaming answer plus TTS"),
    ]
    sw, sh = Inches(2.15), Inches(1.4)
    gap = Inches(0.18)
    total_w = 5 * sw + 4 * Inches(0.28) + 4 * gap
    sx0 = (W - total_w) / 2
    sy = Inches(1.9)
    for i, (num, title, body) in enumerate(steps):
        x = sx0 + i * (sw + Inches(0.28) + gap)
        flow_step(s, num, title, body, x, sy, sw, sh)
        if i < 4:
            arrow(s, x + sw + gap, sy + Inches(0.5))

    cw, ch = Inches(3.85), Inches(1.35)
    cards_y = Inches(3.65)
    card(s, "Value",
         "Guests get immediate service without waiting for reception.",
         Inches(0.5), cards_y, cw, ch, bar_color=GREEN_BAR)
    card(s, "Hotel benefit",
         "Staff handle exceptions while AI handles routine volume.",
         Inches(4.6), cards_y, cw, ch, bar_color=YELLOW_BAR)
    card(s, "Proof in repo",
         "Booking, dining, RAG, guest auth, admin, and telephony routes are implemented.",
         Inches(8.7), cards_y, cw, ch, bar_color=BLUE_BAR)
    footer(s, 5)


def slide06_innovation(prs):
    s = blank_slide(prs)
    bg(s)
    label(s, "Innovation & Uniqueness", Inches(0.5), Inches(0.35))
    heading(s, "StayNep moves from 'chatbot' to autonomous receptionist.",
            Inches(0.5), Inches(0.65), Inches(12.3), size=Pt(26))
    subtext(s, "The Blue Ocean is voice + multilingual hospitality context + live operations.",
            Inches(0.5), Inches(1.35), Inches(11))

    cw, ch = Inches(3.9), Inches(2.8)
    cy = Inches(1.8)

    card(s, "Me Too: basic chatbot",
         "Text-only FAQ\nScripted answers\nNo inventory awareness\nNo booking transaction\nWeak language switching",
         Inches(0.5), cy, cw, ch, bar_color=PINK_BAR, title_size=Pt(14))
    card(s, "Traditional desk",
         "Human warmth\nLimited hours\nSingle-location capacity\nHigh recurring cost\nPeak-hour overload",
         Inches(4.65), cy, cw, ch, bar_color=YELLOW_BAR, title_size=Pt(14))
    card(s, "Only Me: StayNep",
         "Voice-to-voice service\n34-language support\nHotel-specific RAG\nAutonomous booking/dining\nFYI vs escalation routing",
         Inches(8.8), cy, cw, ch, bar_color=GREEN_BAR, title_size=Pt(14))

    txt(s, "Competitive advantage: the assistant connects natural conversation to real hotel\noperations, so it can complete service rather than only describe service.",
        Inches(0.5), Inches(4.95), Inches(12.3), Inches(0.7),
        size=Pt(11.5), color=MUTED, italic=True)
    footer(s, 6)


def slide07_features(prs):
    s = blank_slide(prs)
    bg(s)
    label(s, "Features Overview", Inches(0.5), Inches(0.35))
    heading(s, "The product system covers guest, staff, and admin workflows.",
            Inches(0.5), Inches(0.65), Inches(12.3), size=Pt(26))
    subtext(s, "Major capabilities visible in the application and codebase.",
            Inches(0.5), Inches(1.35), Inches(10))

    cw, ch = Inches(2.95), Inches(1.25)
    row1_y, row2_y = Inches(1.8), Inches(3.2)
    bars1 = [GREEN_BAR, BLUE_BAR, TEAL, YELLOW_BAR]
    bars2 = [PINK_BAR, GREEN_BAR, BLUE_BAR, YELLOW_BAR]

    r1 = [
        ("Voice-to-voice chat", "STT, streaming response, TTS"),
        ("Autonomous booking", "create, modify, cancel"),
        ("Live inventory", "availability and alternatives"),
        ("Dining reservations", "venue, time, party size"),
    ]
    r2 = [
        ("Guest accounts", "My Stay panel and history"),
        ("Admin settings", "hotel config, FAQ, inventory"),
        ("Staff notifications", "FYI emails and urgent tickets"),
        ("Embed/telephony", "web widget, Telnyx, WhatsApp"),
    ]
    for i, (title, body) in enumerate(r1):
        x = Inches(0.5) + i * (cw + Inches(0.18))
        card(s, title, body, x, row1_y, cw, ch, bar_color=bars1[i])
    for i, (title, body) in enumerate(r2):
        x = Inches(0.5) + i * (cw + Inches(0.18))
        card(s, title, body, x, row2_y, cw, ch, bar_color=bars2[i])
    footer(s, 7)


def slide08_ai(prs):
    s = blank_slide(prs)
    bg(s)
    label(s, "AI Integration", Inches(0.5), Inches(0.35))
    heading(s, "AI is the reasoning layer that turns speech into service.",
            Inches(0.5), Inches(0.65), Inches(12.3), size=Pt(26))
    subtext(s, "Without AI, the product becomes a static form plus FAQ.",
            Inches(0.5), Inches(1.35), Inches(10))

    steps = [
        ("A", "Transcribe", "Speech becomes validated text"),
        ("B", "Understand", "Intent, dates, room type, language"),
        ("C", "Retrieve", "Hotel facts and relevant knowledge chunks"),
        ("D", "Decide", "Book, answer, ask, or escalate"),
        ("E", "Speak", "Natural response in guest language"),
    ]
    sw, sh = Inches(2.15), Inches(1.4)
    sx0 = Inches(0.5)
    sy = Inches(1.9)
    gap = Inches(0.2)
    for i, (num, title, body) in enumerate(steps):
        x = sx0 + i * (sw + Inches(0.25) + gap)
        flow_step(s, num, title, body, x, sy, sw, sh)
        if i < 4:
            arrow(s, x + sw + gap, sy + Inches(0.5))

    card(s, "AI functionality",
         "Gemini/OpenAI response engine, RAG augmentation, multilingual replies, booking-state reasoning, and escalation detection.",
         Inches(0.5), Inches(3.6), Inches(5.9), Inches(1.55), bar_color=GREEN_BAR)
    card(s, "If AI is removed",
         "The system loses natural language understanding, voice conversation, multilingual adaptability, and autonomous service decisions.",
         Inches(6.7), Inches(3.6), Inches(5.9), Inches(1.55), bar_color=PINK_BAR)
    footer(s, 8)


def slide09_stack(prs):
    s = blank_slide(prs)
    bg(s)
    label(s, "Technology Stack", Inches(0.5), Inches(0.35))
    heading(s, "The architecture is modern, typed, and deployment-ready.",
            Inches(0.5), Inches(0.65), Inches(12.3), size=Pt(26))
    subtext(s, "Core technologies used by the current repository.",
            Inches(0.5), Inches(1.35), Inches(10))

    cw, ch = Inches(3.9), Inches(1.9)
    row1_y, row2_y = Inches(1.8), Inches(3.9)
    bars = [GREEN_BAR, BLUE_BAR, YELLOW_BAR, GREEN_BAR, BLUE_BAR, YELLOW_BAR]

    items = [
        ("Frontend", "Next.js 16\nReact 19\nTypeScript\nTailwind CSS v4"),
        ("Backend", "App Router API routes\nSSE streaming\nAuth, rate limits\nValidation"),
        ("Database", "Prisma ORM\nPostgreSQL / Prisma Postgres\nHotel, booking, guest models"),
        ("AI tools", "Google Gemini\nOpenAI fallback\nRAG embeddings\nPrompt engine"),
        ("Voice/APIs", "Nemotron/Whisper/Gemini STT\nEdge/Nemotron/OpenAI TTS\nTelnyx, WhatsApp, TingTing"),
        ("Hosting", "Vercel\nEdge-ready routes\nService health checks"),
    ]
    for i, (title, body) in enumerate(items[:3]):
        x = Inches(0.5) + i * (cw + Inches(0.18))
        card(s, title, body, x, row1_y, cw, ch, bar_color=bars[i])
    for i, (title, body) in enumerate(items[3:]):
        x = Inches(0.5) + i * (cw + Inches(0.18))
        card(s, title, body, x, row2_y, cw, ch, bar_color=bars[i + 3])
    footer(s, 9)


def slide10_business(prs):
    s = blank_slide(prs)
    bg(s)
    label(s, "Business Model", Inches(0.5), Inches(0.35))
    heading(s, "A subscription model converts hotel automation into recurring revenue.",
            Inches(0.5), Inches(0.65), Inches(12.3), size=Pt(25))
    subtext(s, "Pricing can scale from small guesthouses to hotel groups.",
            Inches(0.5), Inches(1.35), Inches(10))

    cw, ch = Inches(3.9), Inches(3.0)
    cy = Inches(1.8)

    # Starter
    bx = box(s, Inches(0.5), cy, cw, ch, fill_color=BG_CARD)
    txt(s, "Starter", Inches(0.65), cy + Inches(0.2), cw, Inches(0.4),
        size=Pt(14), bold=True, color=WHITE)
    txt(s, "$49/mo", Inches(0.65), cy + Inches(0.6), cw, Inches(0.6),
        size=Pt(30), bold=False, color=TEAL)
    txt(s, "Includes\n• Small hotels\n• 1–4 languages\n• Core FAQ + voice\n• Email support",
        Inches(0.65), cy + Inches(1.2), cw - Inches(0.2), ch - Inches(1.25),
        size=Pt(11), color=MUTED)

    # Professional
    bx2 = box(s, Inches(4.65), cy, cw, ch, fill_color=BG_CARD)
    txt(s, "Professional", Inches(4.8), cy + Inches(0.2), cw, Inches(0.4),
        size=Pt(14), bold=True, color=WHITE)
    txt(s, "$149/mo", Inches(4.8), cy + Inches(0.6), cw, Inches(0.6),
        size=Pt(30), bold=False, color=YELLOW_BAR)
    txt(s, "Includes\n• Mid-size hotels\n• Booking + dining flows\n• Analytics\n• Priority support",
        Inches(4.8), cy + Inches(1.2), cw - Inches(0.2), ch - Inches(1.25),
        size=Pt(11), color=MUTED)

    # Enterprise
    bx3 = box(s, Inches(8.8), cy, cw, ch, fill_color=BG_CARD)
    txt(s, "Enterprise", Inches(8.95), cy + Inches(0.2), cw, Inches(0.4),
        size=Pt(14), bold=True, color=WHITE)
    txt(s, "Custom", Inches(8.95), cy + Inches(0.6), cw, Inches(0.6),
        size=Pt(30), bold=False, color=BLUE_BAR)
    txt(s, "Includes\n• Hotel chains\n• Custom integrations\n• Telephony volume\n• Dedicated setup",
        Inches(8.95), cy + Inches(1.2), cw - Inches(0.2), ch - Inches(1.25),
        size=Pt(11), color=MUTED)

    txt(s, "Future sustainability: subscriptions, setup fees, PMS/CRM integration packages, and usage-based telephony or message volume.",
        Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.6),
        size=Pt(11), color=MUTED, italic=True)
    footer(s, 10)


def slide11_market(prs):
    s = blank_slide(prs)
    bg(s)
    label(s, "Market & Competitor Analysis", Inches(0.5), Inches(0.35))
    heading(s, "Customers choose the option that combines availability, language, and action.",
            Inches(0.5), Inches(0.65), Inches(12.3), size=Pt(24))
    subtext(s, "Simple comparison table across realistic alternatives.",
            Inches(0.5), Inches(1.35), Inches(10))

    # Table
    tx, ty = Inches(0.5), Inches(1.75)
    tw, th = Inches(12.3), Inches(3.8)
    cols = [Inches(2.1), Inches(2.2), Inches(2.1), Inches(2.4), Inches(2.3)]
    headers = ["Criteria", "Human desk", "FAQ bot", "Generic AI tool", "StayNep"]
    rows_data = [
        ["24/7 service",        "Limited", "Yes",    "Yes",         "Yes"],
        ["Voice support",       "Yes",     "No",     "Partial",     "Yes"],
        ["Hotel-specific data", "Yes",     "Static", "Manual setup","Built in"],
        ["Live booking action", "Manual",  "No",     "No",          "Yes"],
        ["Multilingual scale",  "Limited", "Low",    "Medium",      "34 languages"],
        ["Staff workload",      "High",    "Medium", "Medium",      "Lower"],
    ]

    row_h = Inches(0.48)
    header_h = Inches(0.45)

    # Header row
    hx = tx
    for i, (hdr, cw) in enumerate(zip(headers, cols)):
        fill = BG_CARD if i < 4 else RGBColor(0x0D, 0x2A, 0x22)
        box(s, hx, ty, cw, header_h, fill_color=fill)
        color = TEAL if i == 4 else WHITE
        txt(s, hdr, hx + Inches(0.08), ty + Inches(0.1), cw, header_h,
            size=Pt(10), bold=True, color=color)
        hx += cw

    # Data rows
    for ri, row in enumerate(rows_data):
        ry = ty + header_h + ri * row_h
        rx = tx
        for ci, (cell, cw) in enumerate(zip(row, cols)):
            fill = RGBColor(0x0A, 0x1A, 0x22) if ri % 2 == 0 else BG_CARD
            if ci == 4:
                fill = RGBColor(0x0A, 0x22, 0x18)
            box(s, rx, ry, cw, row_h, fill_color=fill)
            color = TEAL if ci == 4 else (MUTED if ci > 0 else WHITE)
            txt(s, cell, rx + Inches(0.08), ry + Inches(0.1), cw, row_h,
                size=Pt(10), color=color)
            rx += cw

    card(s, "Market opportunity",
         "Nepal's hotels and SMEs need affordable automation that feels local, multilingual, and operationally useful.",
         Inches(6.7), Inches(5.7), Inches(5.8), Inches(1.35), bar_color=TEAL)
    footer(s, 11)


def slide12_sdg(prs):
    s = blank_slide(prs)
    bg(s)
    label(s, "SDG Alignment", Inches(0.5), Inches(0.35))
    heading(s, "The project supports economic growth and digital innovation.",
            Inches(0.5), Inches(0.65), Inches(12.3), size=Pt(26))
    subtext(s, "StayNep aligns most directly with SDG 8 and SDG 9.",
            Inches(0.5), Inches(1.35), Inches(10))

    cw, ch = Inches(5.9), Inches(2.9)
    cy = Inches(1.8)
    card(s, "SDG 8: Decent Work & Economic Growth",
         "Reduces routine operational pressure.\nHelps small hotels serve guests continuously.\nCreates room for staff to focus on higher-value hospitality.\nSupports tourism-driven business growth.",
         Inches(0.5), cy, cw, ch, bar_color=GREEN_BAR, title_size=Pt(13))
    card(s, "SDG 9: Industry, Innovation & Infrastructure",
         "Applies generative AI to hospitality infrastructure.\nMakes advanced automation accessible to SMEs.\nEncourages local technology adoption.\nBuilds scalable digital service channels.",
         Inches(6.7), cy, cw, ch, bar_color=BLUE_BAR, title_size=Pt(13))

    txt(s, "Impact: better guest service, stronger small-hotel competitiveness, and more efficient use of human staff.",
        Inches(0.5), Inches(5.05), Inches(12.3), Inches(0.6),
        size=Pt(12), color=MUTED, italic=True)
    footer(s, 12)


def slide13_journey(prs):
    s = blank_slide(prs)
    bg(s)
    label(s, "Development Journey", Inches(0.5), Inches(0.35))
    heading(s, "The project moved from concept to working product through staged delivery.",
            Inches(0.5), Inches(0.65), Inches(12.3), size=Pt(24))
    subtext(s, "Timeline summary for the Innovation Fest build.",
            Inches(0.5), Inches(1.35), Inches(10))

    timeline_y = Inches(2.6)
    phases = [
        ("Week 1-2",  "Problem\ndiscovery",     TEAL),
        ("Week 3-4",  "Research and\nstack",     BLUE_BAR),
        ("Week 5-6",  "UI and app\nshell",       YELLOW_BAR),
        ("Week 7-9",  "Booking +\nadmin",        TEAL),
        ("Week 10-11","AI + voice\nstack",       BLUE_BAR),
        ("Week 12",   "Testing and\nfixes",      YELLOW_BAR),
        ("Week 13",   "Deployment\nprep",        TEAL),
    ]
    n = len(phases)
    total_w = Inches(12.3)
    step_w = total_w / n
    sx0 = Inches(0.5)

    # Line
    line = slide13_journey.__dict__.get  # just use box
    lbox = box(s, sx0, timeline_y + Inches(0.18), total_w, Pt(2), fill_color=DIVIDER)

    for i, (week, label_text, color) in enumerate(phases):
        cx = sx0 + i * step_w + step_w / 2 - Inches(0.12)
        # dot
        dot = box(s, cx, timeline_y + Inches(0.06), Inches(0.24), Inches(0.24),
                  fill_color=color)
        # week label above
        txt(s, week, cx - Inches(0.4), timeline_y - Inches(0.35),
            Inches(1.1), Inches(0.3), size=Pt(9), bold=True, color=color,
            align=PP_ALIGN.CENTER)
        # phase label below
        txt(s, label_text, cx - Inches(0.4), timeline_y + Inches(0.4),
            Inches(1.1), Inches(0.55), size=Pt(10), color=MUTED,
            align=PP_ALIGN.CENTER)

    card(s, "Milestones achieved",
         "Working voice assistant, autonomous booking flows, dining reservations, RAG knowledge, admin settings, analytics/support screens, guest auth, service health, and deployment configuration.",
         Inches(0.5), Inches(3.75), Inches(11.8), Inches(1.4), bar_color=GREEN_BAR)
    footer(s, 13)


def slide14_challenges(prs):
    s = blank_slide(prs)
    bg(s)
    label(s, "Challenges & Lessons Learned", Inches(0.5), Inches(0.35))
    heading(s, "The hard parts were reliability, context, and production behavior.",
            Inches(0.5), Inches(0.65), Inches(12.3), size=Pt(25))
    subtext(s, "Each challenge led to an implementation improvement.",
            Inches(0.5), Inches(1.35), Inches(10))

    cw, ch = Inches(5.9), Inches(1.3)
    bars = [PINK_BAR, YELLOW_BAR, BLUE_BAR, GREEN_BAR, TEAL, PINK_BAR]
    items = [
        ("Speech accuracy",
         "Added multiple STT options, validation, and silence detection.",
         Inches(0.5), Inches(1.8), PINK_BAR),
        ("Hotel-specific answers",
         "Built config-driven prompts and RAG knowledge chunks.",
         Inches(6.7), Inches(1.8), YELLOW_BAR),
        ("Booking conflicts",
         "Used transactional booking service and alternative suggestions.",
         Inches(0.5), Inches(3.25), BLUE_BAR),
        ("Escalation noise",
         "Separated routine FYI notifications from urgent tickets.",
         Inches(6.7), Inches(3.25), GREEN_BAR),
        ("Multi-tenant context",
         "Introduced slug-based tenant routing and scoped config.",
         Inches(0.5), Inches(4.7), TEAL),
        ("Deployment readiness",
         "Added env examples, health route, tests, and provider fallbacks.",
         Inches(6.7), Inches(4.7), PINK_BAR),
    ]
    for title, body, x, y, bar in items:
        card(s, title, body, x, y, cw, ch, bar_color=bar)
    footer(s, 14)


def slide15_future(prs):
    s = blank_slide(prs)
    bg(s)
    label(s, "Future Enhancements", Inches(0.5), Inches(0.35))
    heading(s, "One more month would deepen integrations and reliability.",
            Inches(0.5), Inches(0.65), Inches(12.3), size=Pt(26))
    subtext(s, "Roadmap from near-term demo polish to scalable product.",
            Inches(0.5), Inches(1.35), Inches(10))

    steps = [
        ("Week 1", "Demo hardening", "Seed data, offline video, QA scripts"),
        ("Week 2", "WhatsApp + telephony", "Production call flows and logs"),
        ("Week 3", "PMS integration", "Real hotel booking system connector"),
        ("Week 4", "Analytics + learning", "FAQ gap review and CSAT improvements"),
        ("Later",  "Mobile app + CRM", "Guest memory, loyalty, IoT controls"),
    ]
    sw, sh = Inches(2.2), Inches(1.4)
    sx0 = Inches(0.45)
    sy = Inches(1.8)
    gap = Inches(0.12)
    for i, (week, title, body) in enumerate(steps):
        x = sx0 + i * (sw + Inches(0.22) + gap)
        flow_step(s, week, title, body, x, sy, sw, sh)
        if i < 4:
            arrow(s, x + sw + gap, sy + Inches(0.5))

    card(s, "Scaling strategy",
         "Package the assistant as an embeddable widget plus managed onboarding for hotels.",
         Inches(0.5), Inches(3.55), Inches(5.9), Inches(1.4), bar_color=GREEN_BAR)
    card(s, "More AI features",
         "Voice personalisation, better multilingual evaluation, and smarter service recovery.",
         Inches(6.7), Inches(3.55), Inches(5.9), Inches(1.4), bar_color=BLUE_BAR)
    footer(s, 15)


def slide16_demo(prs):
    s = blank_slide(prs)
    bg(s)
    label(s, "Live Demo", Inches(0.5), Inches(0.35))
    heading(s, "The demo should prove the core workflow in under five minutes.",
            Inches(0.5), Inches(0.65), Inches(12.3), size=Pt(25))
    subtext(s, "Show login only if needed; move quickly to voice, booking, and output.",
            Inches(0.5), Inches(1.35), Inches(10))

    steps = [
        ("01", "Open assistant", "Show branded hotel ID"),
        ("02", "Ask availability", "Use voice or text input"),
        ("03", "Confirm booking", "Collect dates, room, name, phone"),
        ("04", "Show result", "Booking card + confirmation"),
        ("05", "Switch language", "Repeat a simple guest question"),
    ]
    sw, sh = Inches(2.1), Inches(1.4)
    sx0 = Inches(0.5)
    sy = Inches(1.8)
    gap = Inches(0.22)
    for i, (num, title, body) in enumerate(steps):
        x = sx0 + i * (sw + Inches(0.2) + gap)
        flow_step(s, num, title, body, x, sy, sw, sh)
        if i < 4:
            arrow(s, x + sw + gap, sy + Inches(0.5))

    cw, ch = Inches(3.85), Inches(1.35)
    cards_y = Inches(3.55)
    card(s, "AI features to show",
         "Voice-to-text, streamed response, TTS, multilingual reply, RAG, hotel facts.",
         Inches(0.5), cards_y, cw, ch, bar_color=GREEN_BAR)
    card(s, "Operational output",
         "Live availability, booking confirmation, staff FYI, guest My Stay panel.",
         Inches(4.6), cards_y, cw, ch, bar_color=BLUE_BAR)
    card(s, "Demo backup",
         "Repeat sample demo steps, stable screen captures, and a short recorded backup.",
         Inches(8.7), cards_y, cw, ch, bar_color=YELLOW_BAR)
    footer(s, 16)


def slide17_impact(prs):
    s = blank_slide(prs)
    bg(s)
    label(s, "Impact & Conclusion", Inches(0.5), Inches(0.35))
    heading(s, "StayNep creates measurable value for hotels and guests.",
            Inches(0.5), Inches(0.65), Inches(12.3), size=Pt(26))
    subtext(s, "The project solves a real service bottleneck with a working technical system.",
            Inches(0.5), Inches(1.35), Inches(11))

    stat(s, "24/7", "continuous guest\nservice", Inches(0.5), Inches(2.1), val_color=TEAL)
    stat(s, "34",   "language coverage", Inches(3.3), Inches(2.1), val_color=BLUE_BAR)
    stat(s, "Live", "inventory and\nbooking", Inches(6.0), Inches(2.1), val_color=YELLOW_BAR)
    stat(s, "FYI",  "not ticket spam", Inches(9.0), Inches(2.1), val_color=PINK_BAR)

    cw, ch = Inches(3.9), Inches(1.45)
    cy = Inches(3.65)
    card(s, "Problem solved",
         "Guests get answers and actions without waiting for reception.",
         Inches(0.5), cy, cw, ch, bar_color=GREEN_BAR)
    card(s, "Value created",
         "Hotels reduce routine workload and recover missed inquiries.",
         Inches(4.6), cy, cw, ch, bar_color=YELLOW_BAR)
    card(s, "Innovation achieved",
         "Voice AI is connected to real hospitality operations.",
         Inches(8.7), cy, cw, ch, bar_color=BLUE_BAR)
    footer(s, 17)


def slide18_pitch(prs):
    s = blank_slide(prs)
    bg(s)
    label(s, "Final Pitch", Inches(0.5), Inches(0.35))
    heading(s, "Every hotel deserves a receptionist that never sleeps.",
            Inches(0.5), Inches(0.65), Inches(12.3), size=Pt(26))
    subtext(s, "Why users should adopt it, and why investors should believe in it.",
            Inches(0.5), Inches(1.35), Inches(11))

    cw, ch = Inches(3.9), Inches(2.5)
    cy = Inches(1.8)
    card(s, "Why users adopt",
         "Faster responses\nService in the guest's language\nBooking without waiting\nConsistent answers from hotel data",
         Inches(0.5), cy, cw, ch, bar_color=GREEN_BAR)
    card(s, "Why hotels buy",
         "Lower staff overload\nMore captured inquiries\nScalable front-desk capacity\nClear escalation path",
         Inches(4.6), cy, cw, ch, bar_color=YELLOW_BAR)
    card(s, "Why investors believe",
         "Subscription revenue\nLarge hospitality SME market\nExpandable integrations\nBuilt product foundation",
         Inches(8.7), cy, cw, ch, bar_color=BLUE_BAR)

    txt(s, "The future of hospitality is not only a person behind a desk. It is a service layer that is always available, multilingual, and connected to the hotel's real operations.",
        Inches(0.5), Inches(4.65), Inches(12.3), Inches(0.75),
        size=Pt(12), color=MUTED, italic=True)

    # Footer with team
    txt(s, "Team Think  |  StayNep AI Voice Receptionist",
        Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
        size=Pt(10), color=MUTED, align=PP_ALIGN.CENTER)
    footer(s, 18)


# ── Main ──────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()
    slide01_title(prs)
    slide02_exec_summary(prs)
    slide03_problem(prs)
    slide04_discovery(prs)
    slide05_solution(prs)
    slide06_innovation(prs)
    slide07_features(prs)
    slide08_ai(prs)
    slide09_stack(prs)
    slide10_business(prs)
    slide11_market(prs)
    slide12_sdg(prs)
    slide13_journey(prs)
    slide14_challenges(prs)
    slide15_future(prs)
    slide16_demo(prs)
    slide17_impact(prs)
    slide18_pitch(prs)

    out = "StayNep_Presentation.pptx"
    prs.save(out)
    print(f"✓  Saved {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    build()
